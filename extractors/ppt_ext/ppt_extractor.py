from pathlib import Path
from pptx import Presentation

#paddle pOCR 구동용 import
from paddleocr import PaddleOCR
from PIL import Image
import numpy as np
import io

#.ppt 용 import
import os
import sys
import re
import json
import olefile

"""
#기본 checker stub
def ppt_extractor(file_path: Path):
    print("ppt extractor 연결 성공")
    print(f"파일명: {file_path.name}")
"""

#ocr 가동
ocr = PaddleOCR(
    use_doc_orientation_classify=False,
    use_doc_unwarping=False,
    use_textline_orientation=False,
    lang='korean',
    engine="paddle")

def ppt_extractor(file_path: Path):

    print("ppt extractor 연결 성공")
    print(f"파일명: {file_path.name}")

    extension = file_path.suffix.lower()

    if (extension == ".pptx") :
        ordered_text_per_slide = extract_text_ordered_by_position(file_path)
    else :
        # FIX: process() returns metadata + content for .ppt files, so use only
        # the extracted slide/text content when building the final result.
        ppt_result = process(file_path)
        ordered_text_per_slide = ppt_result.get("content", {})

    # FIX: append() is a list method. This used to be {}, which caused
    # AttributeError: 'dict' object has no attribute 'append'.
    slide_texts = []
    for slide_number in sorted(ordered_text_per_slide):
        slide_text = ordered_text_per_slide[slide_number]
        slide_texts.append(f"Slide {slide_number}:\n{slide_text}\n\n")

    # FIX: Printing the full extracted text can fail on Windows cp949 consoles
    # when slides contain special Unicode characters. Print only a safe summary.
    print(f"extracted slides: {len(slide_texts)}")

    #return slide_texts
    
    
    #슬라이드의 모든 문자열 긁어오기 후 파일명.txt 문서로 출력
    
    output_dir = Path("./output/text")
    output_dir.mkdir(parents=True, exist_ok=True)

    txt_file_name = output_dir / f"{file_path.stem}.txt"

    with open(txt_file_name, "w", encoding="utf-8") as f:
        for slide_number in sorted(ordered_text_per_slide):
            slide_text = ordered_text_per_slide[slide_number]
            f.write(f"Slide {slide_number}:\n{slide_text}\n\n")

    print(f"TXT 저장 완료: {txt_file_name}")

    return slide_texts

def extract_text_from_shape(shape, slide_items):

    # shape에서 텍스트를 추출해서 slide_items 리스트에 append.
    if hasattr(shape, "text") and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:  # 공백만 있는 텍스트는 제외
                slide_items.append((shape.left, shape.top, text))
    
    #하고 싶은 것: 1. 사진인지 파악 2. OCR 돌려서 텍스트 변환 3. 텍스트 저장
    if shape.shape_type == 13 : #MSO_SHAPE.TYPE.PICTURE
        try:
            image = shape.image
            img = Image.open(io.BytesIO(image.blob))
            img = img.convert("RGB")
            img_np = np.array(img)
        
            output_dir = Path("./output")
            output_dir.mkdir(parents=True, exist_ok=True)
            preds = ocr.predict(img_np)

            # FIX: initialize rec_texts first so OCR images with no detected
            # text do not raise an UnboundLocalError.
            rec_texts = []
            for pred in preds:
                rec_texts.extend(pred.get('rec_texts', []))

            texcon = ' '.join(rec_texts)

            #rec_texts = [pred['rec_texts'] for pred in preds]            
            #print(rec_texts)

            #검증용 코드
            #for res in preds:
                #res.print()
                #res.save_to_img(save_path=output_dir) ## Save the processed image
                #res.save_to_json(save_path=output_dir) ## Save the current image's structured result in JSON format
                
            slide_items.append((shape.left, shape.top, texcon))
        
        except Exception as e:
                print(f"OCR 실패: {e}")

    # 테이블 처리
    if shape.has_table:
        table = shape.table
        rows_text = []
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = []
                for paragraph in cell.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        cell_text.append(text)
                row_text.append(' '.join(cell_text))  # 셀 내 텍스트를 결합
            rows_text.append(', '.join(row_text))  # 행 내 텍스트를 결합
        table_text = '\n'.join(rows_text)  # 테이블 전체 텍스트를 결합
        slide_items.append((shape.left, shape.top, table_text))

    # 그룹화된 개체 처리
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for grouped_shape in shape.shapes:
            extract_text_from_shape(grouped_shape, slide_items)

def extract_text_ordered_by_position(pptx_file):
    prs = Presentation(pptx_file)

    slide_texts = {}

    for i, slide in enumerate(prs.slides, start=1):
        slide_items = []
        for shape in slide.shapes:
            extract_text_from_shape(shape, slide_items)

        # 좌상단에서 우하단 방향으로 정렬
        slide_items_sorted = sorted(slide_items, key=lambda x: (x[1], x[0]))
        slide_text = "\n".join(item[2] for item in slide_items_sorted)  # 줄바꿈 추가
        slide_texts[i] = slide_text

    return slide_texts


#.ppt 추출용 코드들

TEXT_HEADER_ATOM = "9F0F"
TEXT_BYTES_ATOM = "A80F"
ENCODING = "unicode_escape"

def hexdump(src, length=16) -> str:
    """
    Returns a hexadecimal dump of a binary string
    (adapted from https://github.com/decalage2/oletools/blob/master/oletools/ezhexviewer.py)

    Args:
        :param src: stream source
        :param length: number of bytes per row
    """
    hex_data = ''
    for i in range(0, len(src), length):
        s = src[i : i + length]
        hex_data += ''.join(["%02X" % x for x in s])
    return hex_data


def process(path: str) -> dict:
    """
    Extract textual content from binary .ppt

    Args:
        :param path: path to the PPT
    """

    # Check if path exists
    if not os.path.exists(path):
        print(f"File {path} does not exist")
        return {}

    try:
        parsed_ppt_dict = {}

        ole = olefile.OleFileIO(path)
        meta = ole.get_metadata()
        parsed_ppt_dict["filename"] = path
        parsed_ppt_dict["slides"] = meta.slides
        parsed_ppt_dict["content"] = {}
        stream = ole.openstream('PowerPoint Document').getvalue()

        hex_data = hexdump(stream)

        matches = list(re.finditer(TEXT_HEADER_ATOM, hex_data))
        matches_spans = [match.span() for match in matches]

        text_counter = 0
        for j in range(len(matches_spans)):
            start_index = int(matches_spans[j][1])
        
            # Check if there is a TextBytesAtom
            is_text_bytes_atom = hex_data[start_index+20:start_index+24] == TEXT_BYTES_ATOM
            if not is_text_bytes_atom:
                continue

            start_index += 8
            # text_type = hex_data[start_index:start_index+2]

            start_index += 16
            # Get the bytes indicating the length of the text
            rec_len_bytes = hex_data[start_index:start_index+4]
            rec_len_bytes_couples = list(zip(rec_len_bytes[0::2], rec_len_bytes[1::2]))
            rec_len_bytes = "".join(["".join(couple) for couple in rec_len_bytes_couples[::-1]])
            # Convert from hex to decimal value
            text_length = int(rec_len_bytes, 16)
            text_start_index = start_index + 8
            hex_text = hex_data[text_start_index:text_start_index+text_length*2]
            byte_string = bytes.fromhex(hex_text)
            result = byte_string.decode(ENCODING, errors="ignore")

            parsed_ppt_dict["content"][str(text_counter)] = result
            text_counter += 1
        
        return parsed_ppt_dict

    except Exception as ex:
        print(f"Something went wrong: {type(ex).__name__}")
        return {}
